"""Populate TennisVision_Presentation.pptx from the course template.

Operates on the copy (the original template is untouched). Each content slide
gets its bullet points in a SINGLE text box (one bulleted list), the other
placeholder bodies are emptied, and the example pictures are swapped for our
figures (pipeline diagram, training curves, dataset sample, annotated frames).
Section tags (TextBox 15) already carry the right names, so they are kept.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from PIL import Image

DECK = "exam/B - Presentation/TennisVision_Presentation.pptx"
A = "exam/B - Presentation/assets"
FIG = "exam/A - Relation/src/figures"
SEL = "output/frames/selected"


def _base_font(shape):
    p = shape.text_frame.paragraphs[0]
    if not p.runs:
        p.add_run()
    b = p.runs[0]
    col = None
    try:
        col = b.font.color.rgb
    except Exception:
        pass
    return b.font.name, col


def set_text(shape, text):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if not p.runs:
        p.add_run()
    p.runs[0].text = text
    for r in p.runs[1:]:
        r._r.getparent().remove(r._r)
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)


def clear_text(shape):
    set_text(shape, "")


def set_bullets(shape, lines, size=18.5):
    """Put all lines as one bulleted list inside a single text box."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    name, col = _base_font(shape)

    def style(run):
        run.font.size = Pt(size)
        if name:
            run.font.name = name
        if col is not None:
            run.font.color.rgb = col

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    p0.runs[0].text = "•  " + lines[0]
    style(p0.runs[0])
    for r in p0.runs[1:]:
        r._r.getparent().remove(r._r)
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    for line in lines[1:]:
        para = tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(6)
        run = para.add_run()
        run.text = "•  " + line
        style(run)


def resize(shape, l, t, w, h):
    shape.left, shape.top = Inches(l), Inches(t)
    shape.width, shape.height = Inches(w), Inches(h)


def _fit(path, l, t, w, h):
    iw, ih = Image.open(path).size
    ar, bar = iw / ih, w / h
    if ar > bar:
        nw, nh = w, w / ar
    else:
        nh, nw = h, h * ar
    return l + (w - nw) / 2, t + (h - nh) / 2, nw, nh


def replace_pic(slide, shape, path):
    l, t, w, h = (shape.left / 914400, shape.top / 914400,
                  shape.width / 914400, shape.height / 914400)
    shape._element.getparent().remove(shape._element)
    fl, ft, fw, fh = _fit(path, l, t, w, h)
    slide.shapes.add_picture(path, Inches(fl), Inches(ft), Inches(fw), Inches(fh))


def add_pic(slide, path, l, t, w, h):
    fl, ft, fw, fh = _fit(path, l, t, w, h)
    slide.shapes.add_picture(path, Inches(fl), Inches(ft), Inches(fw), Inches(fh))


def label(slide, text, l, t, w):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(12)
    r.font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER


prs = Presentation(DECK)
S = [list(s.shapes) for s in prs.slides]

# ---- slide 1: cover ----
set_text(S[0][0], "TennisVision: Player, Ball and Court Analysis "
                  "in Broadcast Tennis Videos with YOLO26")
set_text(S[0][6], "Francesco Zompanti (zompanti.2012601@studenti.uniroma1.it)")
set_text(S[0][7], "Rosario Spaziante (spaziante.2136455@studenti.uniroma1.it)")
set_text(S[0][8], "Luca Panetta (panetta.2136770@studenti.uniroma1.it)")
clear_text(S[0][9])

# ---- slide 3: Context (left image) ----
set_text(S[2][4], "Match Statistics from a Single Broadcast Camera")
set_text(S[2][8], "Why automatic tennis analytics?")
resize(S[2][7], 6.9, 2.6, 8.6, 5.2)
set_bullets(S[2][7], [
    "Shot speed, rally length and player movement guide modern coaching and "
    "broadcasting.",
    "Professional systems (e.g. Hawk-Eye) need several calibrated high-speed "
    "cameras — costly and inaccessible to amateurs.",
    "Our goal: extract these statistics from a single ordinary broadcast "
    "video, with no calibration and no manual annotation.",
])
clear_text(S[2][9]); clear_text(S[2][10])
replace_pic(prs.slides[2], S[2][6], f"{A}/raw_broadcast.png")

# ---- slide 5: Problem & Motivation 1 (left image) ----
set_text(S[4][4], "Challenges of Single-Camera Broadcast Analysis")
set_text(S[4][8], "Why is this hard?")
resize(S[4][7], 6.9, 2.6, 8.6, 5.2)
set_bullets(S[4][7], [
    "The ball is small, fast and motion-blurred, and disappears for several "
    "frames.",
    "Perspective distortion: pixel displacements do not map to physical "
    "speeds.",
    "Broadcast distractors: ball kids, line judges, on-screen graphics.",
])
clear_text(S[4][9]); clear_text(S[4][10])
replace_pic(prs.slides[4], S[4][6], f"{A}/raw_broadcast.png")

# ---- slide 6: Problem & Motivation 2 (full text) ----
set_text(S[5][4], "From Pixels to a Metric Court Space")
set_text(S[5][7], "The core difficulty")
resize(S[5][6], 0.5, 2.6, 15.0, 4.8)
set_bullets(S[5][6], [
    "No calibration and no manual annotation are available at inference time.",
    "Speeds and distances must be physically meaningful, not in pixels.",
    "Detectors pretrained on generic data fail on the tiny tennis ball.",
    "We need court geometry to turn the image into a metric space.",
])
clear_text(S[5][8]); clear_text(S[5][9]); clear_text(S[5][10])

# ---- slide 8: Related Work 1 (left image) ----
set_text(S[7][4], "Detection, Tracking and Court Registration")
set_text(S[7][8], "Related work")
resize(S[7][7], 6.9, 2.6, 8.6, 5.2)
set_bullets(S[7][7], [
    "YOLO family: fast single-stage detectors; we build on YOLO26 "
    "(Ultralytics).",
    "TrackNet: heatmap network for tracking small, fast balls in sport video.",
    "Court registration: classical line/Hough heuristics vs learned keypoint "
    "regression.",
])
clear_text(S[7][9]); clear_text(S[7][10])
replace_pic(prs.slides[7], S[7][6], f"{SEL}/sinner_hardcourt_sec007.png")

# ---- slide 9: Related Work 2 / Gap (full text) ----
set_text(S[8][4], "Our Gap with Respect to the State of the Art")
set_text(S[8][7], "What we do differently")
resize(S[8][6], 0.5, 2.6, 15.0, 4.8)
set_bullets(S[8][6], [
    "Existing open implementations solve only one sub-problem (the ball, or "
    "the court).",
    "Classical court heuristics are fragile to surface colour and lighting.",
    "We fine-tune the ball detector and learn court keypoints instead of "
    "heuristics.",
    "Key difference: a fully metric formulation — every statistic is computed "
    "in real-world meters.",
])
clear_text(S[8][8]); clear_text(S[8][9]); clear_text(S[8][10])

# ---- slide 11: Proposed 1 (big picture = pipeline diagram) ----
set_text(S[10][4], "The TennisVision Pipeline")
replace_pic(prs.slides[10], S[10][7], f"{FIG}/02-PipelineDiagram.png")
resize(S[10][6], 1.3, 6.8, 13.4, 2.0)
set_bullets(S[10][6], [
    "Stages: court keypoints → homography (pixels ↔ meters) → players → ball "
    "→ parabolic smoothing → shot / bounce detection → analytics.",
    "Three YOLO26 models feed one metric court space; every statistic is "
    "computed in real-world meters.",
])
clear_text(S[10][8]); clear_text(S[10][9])

# ---- slide 12: Proposed 2 (full text) ----
set_text(S[11][4], "Geometry, Tracking and Ball Smoothing")
set_text(S[11][7], "How it works")
resize(S[11][6], 0.5, 2.6, 15.0, 4.8)
set_bullets(S[11][6], [
    "Court: YOLO26-pose regresses 14 keypoints; homography H (RANSAC) maps "
    "pixels ↔ meters.",
    "Players: pretrained YOLO26x @1280 px + BoT-SORT; one player per court "
    "half by position, kept across ID changes via spatial continuity.",
    "Ball: fine-tuned YOLO26; off-court and static hotspots removed, then a "
    "piecewise parabolic fit gates outliers and bridges missed detections.",
    "Per-frame homography from a temporal-median keypoint window: follows "
    "camera pan / zoom without reintroducing per-frame jitter.",
])
clear_text(S[11][8]); clear_text(S[11][9]); clear_text(S[11][10])

# ---- slide 13: Proposed 3 (analytics + annotated frame) ----
set_text(S[12][4], "Metric Analytics and Annotated Output")
set_text(S[12][7], "Shots, bounces and live statistics")
resize(S[12][6], 0.5, 2.6, 15.0, 2.0)
set_bullets(S[12][6], [
    "Shots: velocity reversal along the court length; bounces: vertical "
    "reversal in image space.",
    "Per player: shot count, speed, type (serve / volley / groundstroke), and "
    "player + opponent movement speed.",
])
clear_text(S[12][8]); clear_text(S[12][9]); clear_text(S[12][10])
add_pic(prs.slides[12], f"{SEL}/sinner_hardcourt_sec007.png", 4.0, 4.9, 8.0, 3.3)

# ---- slide 15: Dataset 1 (left image) ----
set_text(S[14][4], "Datasets")
set_text(S[14][8], "Two custom-prepared datasets")
resize(S[14][7], 6.9, 2.6, 8.6, 5.2)
set_bullets(S[14][7], [
    "Ball: 578 broadcast frames, single class (Roboflow); 428 / 100 / 50 "
    "train/val/test.",
    "Court: TennisCourtDetector, 8841 frames, 14 keypoints; 6630 train / 2211 "
    "held-out val.",
    "Court JSON converted to Ultralytics pose format; keypoints remapped to a "
    "metric court model.",
])
clear_text(S[14][9]); clear_text(S[14][10])
# raw (non-annotated) samples, one per dataset, stacked on the left
S[14][6]._element.getparent().remove(S[14][6]._element)
label(prs.slides[14], "Court keypoint dataset — raw frame", 0.5, 2.55, 5.9)
add_pic(prs.slides[14], f"{A}/court_raw_hard.png", 0.5, 2.9, 5.9, 2.45)
label(prs.slides[14], "Ball dataset — raw frame", 0.5, 5.5, 5.9)
add_pic(prs.slides[14], f"{A}/ball_raw_grass.jpg", 0.5, 5.85, 5.9, 2.45)

# ---- slide 16: Dataset 2 / protocol (left image = ball curves) ----
set_text(S[15][4], "Training and Evaluation Protocol")
set_text(S[15][8], "Setup")
resize(S[15][7], 6.9, 2.6, 8.6, 5.2)
set_bullets(S[15][7], [
    "Ultralytics / PyTorch from COCO-pretrained YOLO26; ball on a GTX 1650 "
    "(batch 4), court on a Colab T4 (batch 16).",
    "Metrics: mAP@50, mAP@50–95, precision / recall, and ms per image.",
    "Baseline: the COCO zero-shot ‘sports ball’ class vs our fine-tuned "
    "detectors.",
])
clear_text(S[15][9]); clear_text(S[15][10])
replace_pic(prs.slides[15], S[15][6], f"{FIG}/03-BallTrainingCurves.png")

# ---- slide 17: Dataset 3 / Results (one box + table) ----
set_text(S[16][4], "Results")
set_text(S[16][7], "Detection results")
resize(S[16][6], 0.5, 2.6, 15.0, 1.3)
set_bullets(S[16][6], [
    "Fine-tuning lifts ball mAP@50 from ≤0.20 (zero-shot) to 0.905 (YOLO26s).",
    "Court keypoints: 0.994 pose mAP@50 on 2211 held-out frames (5.2 ms/img).",
])
clear_text(S[16][8]); clear_text(S[16][9]); clear_text(S[16][10])
rows = [
    ("Method", "mAP@50", "Recall", "Time"),
    ("YOLO26n zero-shot", "0.064", "0.160", "–"),
    ("YOLO26s zero-shot", "0.200", "0.320", "–"),
    ("YOLO26n fine-tuned", "0.860", "0.805", "12.3 ms"),
    ("YOLO26s fine-tuned", "0.905", "0.880", "17.8 ms"),
]
table = prs.slides[16].shapes.add_table(
    len(rows), 4, Inches(3.2), Inches(4.2), Inches(9.6), Inches(2.4)).table
for ci, w in enumerate((3.6, 2.0, 2.0, 2.0)):
    table.columns[ci].width = Inches(w)
for ri, row in enumerate(rows):
    for cj, val in enumerate(row):
        cell = table.cell(ri, cj)
        cell.text = val
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(14)
        run.font.bold = (ri == 0)

# ---- slide 19: Critical Discussion (3 domain-shift frames) ----
set_text(S[18][4], "Domain Shift Across Surfaces and Lighting")
set_text(S[18][7], "Main limitation: the system is tuned to hardcourt with "
                   "good lighting and degrades on clay and grass.")
clear_text(S[18][6]); clear_text(S[18][8]); clear_text(S[18][10])
frames = [
    (f"{SEL}/sinner_hardcourt_sec007.png", "Hardcourt — works"),
    (f"{SEL}/murray_clay_sec002.png", "Clay — overlay drifts"),
    (f"{SEL}/iga_grass_sec008.png", "Grass — P1 lost"),
]
for (path, cap), x in zip(frames, [0.7, 5.65, 10.6]):
    label(prs.slides[18], cap, x, 2.45, 4.75)
    add_pic(prs.slides[18], path, x, 2.8, 4.75, 2.2)
resize(S[18][9], 0.5, 5.2, 15.0, 2.6)
set_bullets(S[18][9], [
    "Clay: the court homography is unstable, so the metric overlay drifts off "
    "the lines even when ball and players are still detected. Grass: only the "
    "far player (P2) is tracked; the near player drops out.",
    "Training-free fixes tried (CLAHE, lower keypoint threshold, 1280 px + "
    "court-half crops, chamfer-ICP line refinement) only partially help — "
    "root cause is dataset bias.",
])

# ---- slide 21: Conclusion & Future Work (full text) ----
set_text(S[20][4], "Conclusion and Future Work")
set_text(S[20][7], "Summary")
resize(S[20][6], 0.5, 2.6, 15.0, 4.8)
set_bullets(S[20][6], [
    "A complete metric pipeline: custom court-pose and ball detectors, a "
    "pretrained tracker, all analytics in meters.",
    "YOLO26s ball detector: 0.905 mAP@50 / 0.880 recall; court model 0.994 on "
    "2211 held-out frames.",
    "Main limitation: generalization beyond hardcourt and good lighting "
    "(clay / grass).",
    "Future work: ball height for 3D speeds, learned event rules, and in/out "
    "calls from bounces.",
])
clear_text(S[20][8]); clear_text(S[20][9]); clear_text(S[20][10])

prs.save(DECK)
print("saved", DECK)
